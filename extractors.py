#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
extractors.py – Fájltípusonkénti szövegkinyerő modul a Doku RAG rendszerhez.
@author: zsolt

Minden nyilvános függvény szignatúrája:
    extract_<típus>(path: Path) -> str

A főbelépési pont az `extract_text(file_path)`, amely a kiterjesztés alapján
a megfelelő kinyerőt hívja. Üres string-et ad vissza, ha a fájl nem olvasható
vagy a szükséges könyvtár nincs telepítve.
"""

import json
import logging
import re
from pathlib import Path

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Belépési pont
# ---------------------------------------------------------------------------

def extract_text(file_path: Path) -> str:
    """
    Szöveg kinyerése a fájl kiterjesztése alapján.

    Args:
        file_path (Path): A feldolgozandó fájl útvonala.

    Returns:
        str: A kinyert szöveg, vagy üres string olvasási hiba esetén.
    """
    ext = file_path.suffix.lower().lstrip(".")
    extractor = _EXTRACTOR_MAP.get(ext, extract_text_plain)
    try:
        return extractor(file_path)
    except Exception as e:
        logger.warning("Nem sikerült beolvasni: %s – %s", file_path, e)
        return ""


# ---------------------------------------------------------------------------
# Segédfüggvény: szöveges fájl olvasás több encoding-gal
# ---------------------------------------------------------------------------

def _read_raw(path: Path) -> str:
    for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1250"):
        try:
            return path.read_text(encoding=enc)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="replace")


# ---------------------------------------------------------------------------
# Kinyerők
# ---------------------------------------------------------------------------

def extract_text_plain(path: Path) -> str:
    """TXT, MD, HTML, HTM és ismeretlen szöveges fájlok."""
    return _read_raw(path)


def extract_pdf(path: Path) -> str:
    """PDF – pypdf könyvtár."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        parts = []
        for page in reader.pages:
            t = page.extract_text()
            if t:
                parts.append(t)
        return "\n".join(parts)
    except ImportError:
        logger.warning("pypdf nincs telepítve, PDF kihagyva: %s", path)
        return ""


def extract_docx(path: Path) -> str:
    """DOCX – python-docx könyvtár."""
    try:
        from docx import Document
        doc = Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except ImportError:
        logger.warning("python-docx nincs telepítve, DOCX kihagyva: %s", path)
        return ""


def extract_xlsx(path: Path) -> str:
    """XLSX, XLSM – openpyxl könyvtár."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            parts.append(f"[Munkalap: {ws.title}]")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    parts.append("\t".join(cells))
        return "\n".join(parts)
    except ImportError:
        logger.warning("openpyxl nincs telepítve, XLSX kihagyva: %s", path)
        return ""


def extract_xls(path: Path) -> str:
    """XLS (legacy) – pandas + xlrd."""
    try:
        import pandas as pd
        df = pd.read_excel(str(path), engine="xlrd")
        return df.to_string(index=False)
    except ImportError:
        logger.warning("pandas/xlrd nincs telepítve, XLS kihagyva: %s", path)
        return ""


def extract_pptx(path: Path) -> str:
    """PPTX – python-pptx könyvtár."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"[Dia {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        return "\n".join(parts)
    except ImportError:
        logger.warning("python-pptx nincs telepítve, PPTX kihagyva: %s", path)
        return ""


def extract_odt(path: Path) -> str:
    """ODT – odfpy könyvtár."""
    try:
        from odf.opendocument import load as odf_load
        from odf.text import P
        doc = odf_load(str(path))
        parts = []
        for p in doc.getElementsByType(P):
            t = "".join(n.data for n in p.childNodes if n.nodeType == 3)
            if t.strip():
                parts.append(t.strip())
        return "\n".join(parts)
    except ImportError:
        logger.warning("odfpy nincs telepítve, ODT kihagyva: %s", path)
        return ""


def extract_rtf(path: Path) -> str:
    """RTF – striprtf könyvtár, fallback: regex."""
    try:
        import striprtf.striprtf as sr
        raw = path.read_bytes().decode("latin-1", errors="replace")
        return sr.rtf_to_text(raw)
    except ImportError:
        raw = _read_raw(path)
        return re.sub(r"\\[a-z]+\d* ?|\{|\}", " ", raw)


def extract_csv(path: Path) -> str:
    """CSV – pandas, max 5000 sor."""
    try:
        import pandas as pd
        df = pd.read_csv(str(path), nrows=5000)
        return df.to_string(index=False)
    except ImportError:
        return _read_raw(path)


def extract_json(path: Path) -> str:
    """JSON – pretty-print formátumban."""
    try:
        raw = _read_raw(path)
        data = json.loads(raw)
        return json.dumps(data, ensure_ascii=False, indent=2)
    except Exception:
        return _read_raw(path)


def extract_ods(path: Path) -> str:
    """ODS (OpenDocument táblázat) – pandas + odf engine."""
    try:
        import pandas as pd
        df = pd.read_excel(str(path), engine="odf")
        return df.to_string(index=False)
    except ImportError:
        logger.warning("pandas/odfpy nincs telepítve, ODS kihagyva: %s", path)
        return ""


# ---------------------------------------------------------------------------
# Kiterjesztés → kinyerő leképzés
# ---------------------------------------------------------------------------

_EXTRACTOR_MAP: dict = {
    "txt":   extract_text_plain,
    "md":    extract_text_plain,
    "html":  extract_text_plain,
    "htm":   extract_text_plain,
    "pdf":   extract_pdf,
    "docx":  extract_docx,
    "xlsx":  extract_xlsx,
    "xlsm":  extract_xlsx,
    "xls":   extract_xls,
    "pptx":  extract_pptx,
    "odt":   extract_odt,
    "ods":   extract_ods,
    "rtf":   extract_rtf,
    "csv":   extract_csv,
    "json":  extract_json,
}

# Nyilvánosan lekérdezhető támogatott kiterjesztések
SUPPORTED_EXTENSIONS: set[str] = set(_EXTRACTOR_MAP.keys())
